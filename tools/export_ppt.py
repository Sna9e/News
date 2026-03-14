import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import datetime
import tools.chart_generator as cg 

def clear_placeholders(slide):
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape.element
            sp.getparent().remove(sp)

def generate_ppt(data, timeline_data, filename, model_name):
    template_path = "template.pptx"
    if os.path.exists(template_path):
        try: prs = Presentation(template_path)
        except Exception: prs = Presentation()
    else: prs = Presentation()
        
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    clear_placeholders(slide) 
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_box.text_frame.paragraphs[0].text = "高管战报：前沿情报深度分析"
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_box.text_frame.paragraphs[0].text = f"生成日期: {datetime.date.today()}"
    subtitle_box.text_frame.paragraphs[0].font.size = Pt(16)
    subtitle_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    if timeline_data:
        for t_data in timeline_data:
            if not t_data['events']: continue
            chunk_size = 7 
            events = t_data['events']
            for i in range(0, len(events), chunk_size):
                chunk = events[i:i + chunk_size]
                slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
                clear_placeholders(slide) 
                
                t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(0.8))
                t_box.text_frame.paragraphs[0].text = f"⏱️ {t_data['topic']} - 核心时间线"
                t_box.text_frame.paragraphs[0].font.size = Pt(24)
                t_box.text_frame.paragraphs[0].font.bold = True
                
                b_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5))
                tf = b_box.text_frame
                tf.word_wrap = True
                for idx, item in enumerate(chunk):
                    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                    p.text = f"[{item.date}] {item.event} ({item.source})"
                    p.font.size = Pt(14)
                    p.space_after = Pt(8)

    for section in data:
        if not section['data']: continue
        
        finance = section.get('finance', {})
        if finance.get('is_public') and finance.get('chart_path'):
            f_slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
            clear_placeholders(f_slide)
            
            t_box = f_slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
            t_box.text_frame.paragraphs[0].text = f"📊 {section['topic']} ({finance['ticker']}) - 量化面与事件催化"
            t_box.text_frame.paragraphs[0].font.size = Pt(22)
            t_box.text_frame.paragraphs[0].font.bold = True
            
            b_box = f_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.0), Inches(2.8))
            tf = b_box.text_frame
            
            trend_icon = "🔺" if finance['change_pct'] > 0 else "🔻"
            color = RGBColor(192, 0, 0) if finance['change_pct'] > 0 else RGBColor(0, 150, 0)
            p_price = tf.paragraphs[0]
            p_price.text = f"{finance['current_price']} {finance['currency']}  {trend_icon} {finance['change_pct']}%"
            p_price.font.size = Pt(24)
            p_price.font.bold = True
            p_price.font.color.rgb = color
            
            metrics = [
                f"▪ 估值水平: {finance['pe_pb']}",
                f"▪ 股权风险溢价: {finance['erp']}",
                f"▪ 总市值: {finance['market_cap']}"
            ]
            for m in metrics:
                p = tf.add_paragraph()
                p.text = m
                p.font.size = Pt(13)
                p.space_before = Pt(12)

            if os.path.exists(finance['chart_path']):
                f_slide.shapes.add_picture(finance['chart_path'], Inches(4.5), Inches(1.2), width=Inches(5.0))

            cat = finance.get('catalysts', {})
            boxes_data = [
                ("🏛️ 政策与监管", cat.get('policy', "近期无重大政策催化")),
                ("💰 财报与盈利", cat.get('earnings', "未见核心财报数据")),
                ("🚀 产业标志事件", cat.get('landmark', "产业层级平稳")),
                ("🔄 市场风格轮动", cat.get('style', "风格未见明显切换"))
            ]
            
            for i, (title, content) in enumerate(boxes_data):
                x_pos = 0.5 + (i * 2.2) 
                c_box = f_slide.shapes.add_textbox(Inches(x_pos), Inches(4.5), Inches(2.1), Inches(2.5))
                c_tf = c_box.text_frame
                c_tf.word_wrap = True
                
                p_t = c_tf.paragraphs[0]
                p_t.text = title
                p_t.font.size = Pt(12)
                p_t.font.bold = True
                p_t.font.color.rgb = RGBColor(0, 51, 102)
                
                p_c = c_tf.add_paragraph()
                p_c.text = content
                p_c.font.size = Pt(11)
                p_c.space_before = Pt(6)

        # ====================================================
        # 🌟 全新大招：智库多专家会审辩论排版页
        # ====================================================
        committee = section.get('committee')
        if committee:
            c_slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
            clear_placeholders(c_slide)

            # 顶部标题
            t_box = c_slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
            t_box.text_frame.paragraphs[0].text = f"⚖️ {section['topic']} - 科技智库多专家深度会审"
            t_box.text_frame.paragraphs[0].font.size = Pt(24)
            t_box.text_frame.paragraphs[0].font.bold = True

            # 【左侧】绿色阵营：创新战略官 (DeepSeek主笔)
            pro_box = c_slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.3), Inches(3.5))
            pro_tf = pro_box.text_frame
            pro_tf.word_wrap = True
            
            p_pro_title = pro_tf.paragraphs[0]
            p_pro_title.text = "🟢 创新战略官 (突破与潜力)"
            p_pro_title.font.size = Pt(16)
            p_pro_title.font.bold = True
            p_pro_title.font.color.rgb = RGBColor(0, 128, 0)

            for pt in committee.get('pro_points', []):
                p = pro_tf.add_paragraph()
                p.text = f"• {pt}"
                p.font.size = Pt(13)
                p.space_before = Pt(8)

           # 【右侧】红色阵营：产业风控官 (Doubao主笔)
            con_box = c_slide.shapes.add_textbox(Inches(5.0), Inches(1.3), Inches(4.3), Inches(3.5))
            con_tf = con_box.text_frame
            con_tf.word_wrap = True
            
            p_con_title = con_tf.paragraphs[0]
            p_con_title.text = "🔴 产业风控官 (Doubao主笔)"
            p_con_title.font.size = Pt(16)
            p_con_title.font.bold = True
            p_con_title.font.color.rgb = RGBColor(192, 0, 0)

            for pt in committee.get('con_points', []):
                p = con_tf.add_paragraph()
                p.text = f"• {pt}"
                p.font.size = Pt(13)
                p.space_before = Pt(8)

            # 【底部】智库总编的最终决断 (带权重倾向)
            verdict_box = c_slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(8.8), Inches(1.8))
            verdict_tf = verdict_box.text_frame
            verdict_tf.word_wrap = True
            
            p_v_title = verdict_tf.paragraphs[0]
            p_v_title.text = "🏛️ 智库总编最终评估决断："
            p_v_title.font.size = Pt(15)
            p_v_title.font.bold = True
            p_v_title.font.color.rgb = RGBColor(0, 51, 102)

            p_v_content = verdict_tf.add_paragraph()
            p_v_content.text = committee.get('chief_verdict', '')
            p_v_content.font.size = Pt(13)
            p_v_content.line_spacing = 1.1
            p_v_content.space_before = Pt(6)

        # 常规新闻页保持不变
        for news in section['data']:
            slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
            clear_placeholders(slide) 
            
            t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(0.8))
            t_box.text_frame.paragraphs[0].text = news.title
            t_box.text_frame.paragraphs[0].font.size = Pt(22)
            t_box.text_frame.paragraphs[0].font.bold = True

            has_chart = hasattr(news, 'chart_info') and news.chart_info.has_chart and len(news.chart_info.labels) > 0
            text_width = 5.2 if has_chart else 9.0 
            
            b_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(text_width), Inches(5))
            tf = b_box.text_frame
            tf.word_wrap = True
            
            tf.paragraphs[0].text = f"📌 来源: {news.source}  |  🕒 {news.date_check}  |  🔥 热度: {'⭐'*news.importance}"
            tf.paragraphs[0].font.size = Pt(12) 
            tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
            tf.paragraphs[0].space_after = Pt(10)
            
            for line in news.summary.split('\n'):
                line = line.strip()
                if not line: continue
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(13)
                p.space_after = Pt(6) 
                if line.startswith("【"):
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0, 51, 102)

            news_url = getattr(news, 'url', '') 
            if news_url:
                p_link = tf.add_paragraph()
                p_link.text = f"🔗 溯源查证: 点击查看原文"
                p_link.font.size = Pt(11)
                p_link.font.color.rgb = RGBColor(0, 112, 192) 
                p_link.runs[0].hyperlink.address = news_url

            if has_chart:
                try:
                    chart_img = cg.generate_and_download_chart(news.chart_info.chart_title, news.chart_info.labels, news.chart_info.values, news.chart_info.chart_type)
                    if chart_img and os.path.exists(chart_img):
                        slide.shapes.add_picture(chart_img, Inches(5.8), Inches(1.5), width=Inches(3.8))
                except Exception: pass

    path = f"{filename}.pptx"
    prs.save(path)
    return path
